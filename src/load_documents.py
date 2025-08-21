import os
import json
import docx2txt
from pypdf import PdfReader

# 新增的库
try:
    from pptx import Presentation
except ImportError:
    pass

try:
    import pandas as pd
except ImportError:
    pass


def load_text_file(file_path):
    """读取 TXT 文件"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    return {
        "content": content,
        "pages": 1,
        "word_count": len(content.split()),
        "char_count": len(content),
        "type": "txt"
    }


def load_pdf_file(file_path):
    """读取 PDF 文件"""
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"

    return {
        "content": text.strip(),
        "pages": len(reader.pages),
        "word_count": len(text.split()),
        "char_count": len(text),
        "type": "pdf"
    }


def load_docx_file(file_path):
    """读取 Word (.docx) 文件"""
    try:
        content = docx2txt.process(file_path)
        return {
            "content": content.strip(),
            "pages": 1,  # Word 不好统计页数，简化处理
            "word_count": len(content.split()),
            "char_count": len(content),
            "type": "docx"
        }
    except Exception as e:
        print(f"[错误] 无法读取 DOCX 文件 {file_path}: {e}")
        return None


def load_md_file(file_path):
    """读取 Markdown (.md) 文件"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    return {
        "content": content.strip(),
        "pages": 1,
        "word_count": len(content.split()),
        "char_count": len(content),
        "type": "md"
    }


def load_pptx_file(file_path):
    """读取 PowerPoint (.pptx) 文件"""
    try:
        pres = Presentation(file_path)
        text = ""
        slide_count = len(pres.slides)

        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return {
            "content": text.strip(),
            "pages": slide_count,
            "word_count": len(text.split()),
            "char_count": len(text),
            "type": "pptx"
        }
    except Exception as e:
        print(f"[错误] 无法读取 PPTX 文件 {file_path}: {e}")
        return None


def load_xlsx_file(file_path):
    """读取 Excel (.xlsx) 文件（合并所有 sheet）"""
    try:
        df = pd.read_excel(file_path, sheet_name=None)  # 读所有 sheet
        text = ""
        sheet_info = []

        for sheet_name, sheet_df in df.items():
            sheet_text = sheet_df.to_string()
            text += f"\n\n=== {sheet_name} ===\n{sheet_text}\n"
            sheet_info.append(sheet_name)

        return {
            "content": text.strip(),
            "pages": len(df),  # sheet 数量
            "word_count": len(text.split()),
            "char_count": len(text),
            "type": "xlsx",
            "sheets": sheet_info
        }
    except Exception as e:
        print(f"[错误] 无法读取 XLSX 文件 {file_path}: {e}")
        return None


def load_document(file_path):
    """根据文件类型调用对应的加载函数，返回结构化数据"""
    ext = os.path.splitext(file_path)[1].lower()
    filename = os.path.basename(file_path)

    result = {
        "filename": filename,
        "file_path": file_path,
        "status": "success"
    }

    try:
        if ext == ".txt":
            data = load_text_file(file_path)
        elif ext == ".pdf":
            data = load_pdf_file(file_path)
        elif ext == ".docx":
            data = load_docx_file(file_path)
        elif ext == ".md":
            data = load_md_file(file_path)
        elif ext == ".pptx":
            data = load_pptx_file(file_path)
        elif ext == ".xlsx":
            data = load_xlsx_file(file_path)
        else:
            print(f"[警告] 不支持的文件类型: {filename}")
            result["status"] = "unsupported"
            result["content"] = ""
            return result

        if data is None:
            result["status"] = "failed"
            result["content"] = ""
        else:
            result.update(data)  # 合并内容和元数据

    except Exception as e:
        print(f"[异常] 读取文件失败 {filename}: {e}")
        result["status"] = "failed"
        result["content"] = ""

    return result


def load_all_documents(docs_folder):
    """加载 docs/ 文件夹下所有支持的文档"""
    documents = []
    supported_exts = {".txt", ".pdf", ".docx", ".md", ".pptx", ".xlsx"}

    for filename in os.listdir(docs_folder):
        file_path = os.path.join(docs_folder, filename)
        ext = os.path.splitext(filename)[1].lower()

        if os.path.isfile(file_path) and ext in supported_exts:
            print(f"正在读取: {filename}")
            doc_data = load_document(file_path)
            documents.append(doc_data)

    return documents


# === 主程序 ===
if __name__ == "__main__":
    # 动态获取 docs 路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    docs_folder = os.path.join(current_dir, "..", "docs")
    docs_folder = os.path.abspath(docs_folder)

    print(f"🔍 正在查找文档文件夹: {docs_folder}")

    if not os.path.exists(docs_folder):
        print(f"❌ 错误：文件夹不存在！请检查路径")
        print(f"   当前工作目录: {os.getcwd()}")
        print(f"   脚本位置: {__file__}")
        exit(1)

    print(f"✅ 找到文件夹，正在加载文档...")

    docs = load_all_documents(docs_folder)

    # 过滤掉失败的
    successful_docs = [d for d in docs if d["status"] == "success"]

    print("\n" + "=" * 60)
    print("✅ 文档加载完成！")
    print(f"📊 总计文件: {len(docs)} | 成功: {len(successful_docs)} | 失败: {len(docs) - len(successful_docs)}")
    print("=" * 60)

    for doc in successful_docs:
        print(f"\n📄 文件: {doc['filename']} ({doc['type']})")
        print(f"📏 页数: {doc['pages']}, 字数: {doc['word_count']}, 字符数: {doc['char_count']}")
        print(f"📝 内容预览:\n{doc['content'][:300]}...")

    # === 保存为 JSON 文件 ===
    output_file = os.path.join(current_dir, "docs.json")
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(docs, f, ensure_ascii=False, indent=2)

    print(f"\n💾 文档数据已保存至: {output_file}")